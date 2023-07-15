from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
import json
from PIL import Image


def config_read():
    with open("..\\config\\config.json") as settings_json:
        ppt_configs = json.load(settings_json)
    return ppt_configs["presentation"]


def slide_creator(slide_config):
    top = Pt(120)
    title_type = slide_config["type"]
    title_text = slide_config["title"]
    content = slide_config["content"]
    if title_type == "title":
        #to create a title slide
        slide_type = 0
    elif title_type == "list":
        #to create a slide with list
        slide_type = 1
    else:
        #to create a slide with only a title and empty space
        slide_type = 5

    new_slide = pres.slide_layouts[slide_type]
    slide = pres.slides.add_slide(new_slide)
    slide.shapes.title.text = title_text

    if title_type == "title":
        slide.placeholders[1].text = content

    elif title_type == "list":
        list_rows = slide_config["content"]

        for i, list_row in enumerate(list_rows):
            level = list_row["level"]-1
            if i == 0:
                #every member will be added to this line always runs first
                list_first_line = slide.placeholders[1]
                list_first_line.level = level
                list_first_line.text = list_row["text"]
            else:
                list_line = list_first_line.text_frame.add_paragraph()
                list_line.level = level
                list_line.text = list_row["text"]

    elif title_type == "text":
        width, height = Pt(650), Pt(425)
        left = (pres.slide_width - width) / 2
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_box.text = slide_config["content"]
        tf = text_box.text_frame
        tf.word_wrap = True

    elif title_type == "picture":
        image_path = "..\\images\\" + slide_config["content"]
        im = Image.open(image_path)
        picture_width = im.width
        image_left = Pt((pres.slide_width.pt - picture_width) / 2)
        slide.shapes.add_picture(image_path, image_left, top)

    elif title_type == "plot":
        data_path = "..\\data\\"+slide_config["content"]
        chart_data = XyChartData()
        series = chart_data.add_series("")
        data = [i.strip().split(";") for i in open(data_path).readlines()]
        for data_point in data:
            series.add_data_point(data_point[0], data_point[1])
        chart_width, chart_length = Pt(480), Pt(355)
        left = (pres.slide_width - chart_width) / 2
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES, left, top, chart_width, chart_length, chart_data
        ).chart
        chart.has_title = False
        #changing the axis labels
        x_axis = chart.category_axis
        x_axis.axis_title.text_frame.text = slide_config["configuration"]["x-label"]
        y_axis = chart.value_axis
        y_axis.axis_title.text_frame.text = slide_config["configuration"]["y-label"]


if __name__ == "__main__":
    pres_name = "default_name.pptx"
    #pres_name = input("Presentation name: ") + ".pptx"
    pres = Presentation()
    slides_conf = config_read()
    for slide_conf in slides_conf:
        slide_creator(slide_conf)
    pres.save(pres_name)
