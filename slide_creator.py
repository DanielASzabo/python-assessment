from PIL import Image
from pptx.util import Pt
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE


def create_title_slide(presentation, title_text, subtitle_text):
    new_slide = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(new_slide)
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text


def create_list_slide(presentation, title_text, list_data):
    new_slide = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(new_slide)
    slide.shapes.title.text = title_text
    for i, list_row in enumerate(list_data):
        level = list_row["level"] - 1
        if i == 0:
            # every member will be added to this line always runs first
            list_first_line = slide.placeholders[1]
            list_first_line.level = level
            list_first_line.text = list_row["text"]
        else:
            list_line = list_first_line.text_frame.add_paragraph()
            list_line.level = level
            list_line.text = list_row["text"]


def create_text_slide(presentation, title_text, text_box_text):
    new_slide = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(new_slide)
    slide.shapes.title.text = title_text
    width, height, top = Pt(650), Pt(425), Pt(120)
    left = (presentation.slide_width - width) / 2
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_box.text = text_box_text
    tf = text_box.text_frame
    tf.word_wrap = True


def create_picture_slide(presentation, title_text, image_name):
    new_slide = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(new_slide)
    slide.shapes.title.text = title_text
    image_path = "..\\images\\" + image_name
    im = Image.open(image_path)
    picture_width = im.width
    top = Pt(120)
    image_left = Pt((presentation.slide_width.pt - picture_width) / 2)
    slide.shapes.add_picture(image_path, image_left, top)

def create_plot_slide(presentation, title_text, data_name, configuration):
    new_slide = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(new_slide)
    slide.shapes.title.text = title_text
    data_path = "..\\data\\" + data_name
    chart_data = XyChartData()
    series = chart_data.add_series("")
    data = [i.strip().split(";") for i in open(data_path).readlines()]
    for data_point in data:
        series.add_data_point(data_point[0], data_point[1])
    chart_width, chart_length, top = Pt(480), Pt(355), Pt(120)
    left = (presentation.slide_width - chart_width) / 2
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER_LINES, left, top, chart_width, chart_length, chart_data
    ).chart
    chart.has_title = False
    # changing the axis labels
    x_axis = chart.category_axis
    x_axis.axis_title.text_frame.text = configuration["x-label"]
    y_axis = chart.value_axis
    y_axis.axis_title.text_frame.text = configuration["y-label"]